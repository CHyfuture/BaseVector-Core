"""
PPT文档解析器
"""
from pathlib import Path
from typing import Any, Dict

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from ability.operators.parsers.base_parser import BaseParser


class PPTXParser(BaseParser):
    """PPT文档解析器，使用python-pptx"""

    def __init__(self, config: Dict[str, Any] = None):
        """
        初始化PPT解析器

        Args:
            config: 配置字典
        """
        super().__init__(config)
        self.supported_extensions = [".pptx", ".ppt"]

        if Presentation is None:
            raise ImportError(
                "python-pptx is required for PPT parsing. "
                "Install it with: pip install python-pptx"
            )

    def _initialize(self) -> None:
        """初始化PPT解析器"""
        if Presentation is None:
            raise ImportError("python-pptx is not installed")
        self.logger.info("PPT parser initialized")

    def _parse(self, file_path: Path, **kwargs) -> Dict[str, Any]:
        """
        解析PPT文档

        Args:
            file_path: PPT文件路径
            **kwargs: 额外的处理参数

        Returns:
            解析结果字典
        """
        prs = Presentation(str(file_path))

        # 提取元数据
        core_props = prs.core_properties
        metadata = {
            "title": core_props.title or "",
            "author": core_props.author or "",
            "subject": core_props.subject or "",
            "keywords": core_props.keywords or "",
            "comments": core_props.comments or "",
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
            "slide_count": len(prs.slides),
        }

        content_parts = []
        structure = {
            "slides": [],
        }

        # 提取幻灯片内容
        for slide_idx, slide in enumerate(prs.slides):
            slide_content = []
            slide_structure = {
                "slide_number": slide_idx + 1,
                "shapes": [],
            }

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()

                    # 判断是否为标题
                    if shape.is_placeholder:
                        placeholder = shape.placeholder
                        if placeholder and placeholder.placeholder_format.idx == 0:
                            # 标题占位符
                            content_parts.append(f"## Slide {slide_idx + 1}: {text}\n\n")
                            slide_structure["shapes"].append(
                                {
                                    "type": "title",
                                    "text": text,
                                }
                            )
                        else:
                            slide_content.append(text)
                            slide_structure["shapes"].append(
                                {
                                    "type": "text",
                                    "text": text,
                                }
                            )
                    else:
                        slide_content.append(text)
                        slide_structure["shapes"].append(
                            {
                                "type": "text",
                                "text": text,
                            }
                        )

            if slide_content:
                content_parts.append(f"### Slide {slide_idx + 1} Content\n\n")
                content_parts.append("\n\n".join(slide_content))
                content_parts.append("\n\n")

            structure["slides"].append(slide_structure)

        content = "\n".join(content_parts)

        return {
            "content": content,
            "metadata": metadata,
            "structure": structure,
        }
